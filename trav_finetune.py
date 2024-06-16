"""
1. finetune the segmentation model on our traversability dataset
2. infer the finetuned 
"""

from tqdm import tqdm
import network
import utils
import os
import random
import argparse
import numpy as np
import attacks
from pptx import Presentation
from pptx.util import Inches
import pandas as pd

from torch.utils import data
from local_datasets import IndoorTrav
from utils import ext_transforms as et
from metrics import StreamSegMetrics

import torch
import torch.nn as nn
from utils.visualizer import Visualizer

from PIL import Image
import matplotlib
import matplotlib.pyplot as plt

from torch.autograd import Variable


def get_argparser():
    parser = argparse.ArgumentParser()

    # Datset Options
    parser.add_argument("--data_root", type=str, default='/home/qiyuan/2023spring/segmentation_indoor_images',
                        help="path to Dataset")
    parser.add_argument("--dataset", type=str, default='trav',
                        choices=['voc', 'cityscapes', 'trav'], help='Name of dataset')
    parser.add_argument("--scenes", type=list, default=['elb', 'erb', 'uc', 'woh'],
                        choices=['elb', 'erb', 'uc', 'nh', 'woh'], help='Name of dataset')
    parser.add_argument("--num_classes", type=int, default=None,
                        help="num classes (default: None)")

    # Deeplab Options
    available_models = sorted(name for name in network.modeling.__dict__ if name.islower() and \
                              not (name.startswith("__") or name.startswith('_')) and callable(
                              network.modeling.__dict__[name])
                              )
    parser.add_argument("--model", type=str, default='deeplabv3plus_resnet101',
                        choices=available_models, help='model name')
    parser.add_argument("--separable_conv", action='store_true', default=False,
                        help="apply separable conv to decoder and aspp")
    parser.add_argument("--output_stride", type=int, default=16, choices=[8, 16])

    # Train Options
    parser.add_argument("--test_only", action='store_true', default=False)
    parser.add_argument("--save_val_results", action='store_true', default=True,
                        help="save segmentation results to \"./results\"")
    parser.add_argument("--total_itrs", type=int, default=3e4,
                        help="total iter number (default: 3e4)")
    parser.add_argument("--lr", type=float, default=0.01,
                        help="learning rate (default: 0.01)")
    parser.add_argument("--lr_policy", type=str, default='poly', choices=['poly', 'step'],
                        help="learning rate scheduler policy")
    parser.add_argument("--step_size", type=int, default=10000)
    parser.add_argument("--crop_val", action='store_true', default=False,
                        help='crop validation (default: False)')
    parser.add_argument("--batch_size", type=int, default=16,
                        help='batch size (default: 16)')
    parser.add_argument("--val_batch_size", type=int, default=4,
                        help='batch size for validation (default: 4)')
    parser.add_argument("--crop_size", type=int, default=480)

    parser.add_argument("--ckpt", default='checkpoints/best_deeplabv3plus_resnet101_cityscapes_os16.pth.tar',
                        type=str, help="restore from checkpoint")
    parser.add_argument("--continue_training", action='store_true', default=False)
    parser.add_argument("--eps", type=float, default=0.01,
                        help="PGD attack epsilon")
    parser.add_argument("--alpha", type=float, default=100.0,
                        help="PGD attack alpha")
    parser.add_argument("--loss_type", type=str, default='cross_entropy',
                        choices=['cross_entropy', 'focal_loss'], help="loss type (default: False)")
    parser.add_argument("--gpu_id", type=str, default='0,1',
                        help="GPU ID")
    parser.add_argument("--weight_decay", type=float, default=1e-4,
                        help='weight decay (default: 1e-4)')
    parser.add_argument("--random_seed", type=int, default=444,
                        help="random seed (default: 1)")
    parser.add_argument("--print_interval", type=int, default=10,
                        help="print interval of loss (default: 10)")
    parser.add_argument("--val_interval", type=int, default=1e4,
                        help="epoch interval for eval (default: 100)")
    parser.add_argument("--download", action='store_true', default=False,
                        help="download datasets")

    # PASCAL VOC Options
    parser.add_argument("--year", type=str, default='2012',
                        choices=['2012_aug', '2012', '2011', '2009', '2008', '2007'], help='year of VOC')

    # Visdom options
    parser.add_argument("--enable_vis", action='store_true', default=False,
                        help="use visdom for visualization")
    parser.add_argument("--vis_port", type=str, default='13570',
                        help='port for visdom')
    parser.add_argument("--vis_env", type=str, default='main',
                        help='env for visdom')
    parser.add_argument("--vis_num_samples", type=int, default=8,
                        help='number of samples for visualization (default: 8)')
    return parser


def get_dataset(opts):
    """ Dataset And Augmentation
    """
    if opts.dataset == 'trav':
        train_transform = et.ExtCompose([
            # et.ExtResize( 512 ),
            et.ExtRandomCrop(size=(opts.crop_size, opts.crop_size)),
            et.ExtColorJitter(brightness=0.5, contrast=0.5, saturation=0.5),
            et.ExtRandomHorizontalFlip(),
            et.ExtToTensor(),
            et.ExtNormalize(mean=[0.5174, 0.4857, 0.5054],
                            std=[0.2726, 0.2778, 0.2861]),
        ])

        val_transform = et.ExtCompose([
            # et.ExtResize( 512 ),
            et.ExtToTensor(),
            et.ExtNormalize(mean=[0.5174, 0.4857, 0.5054],
                            std=[0.2726, 0.2778, 0.2861]),
        ])

        train_dst = IndoorTrav(opts.data_root, 'train', opts.scenes,
                               transform=train_transform)
        val_dst = IndoorTrav(opts.data_root, 'val', opts.scenes,
                             transform=val_transform)

    return train_dst, val_dst


def validate(opts, model, loader, device, metrics, iter, ret_samples_ids=None):
    """Do validation and return specified samples"""
    metrics.reset()
    ret_samples = []
    if opts.save_val_results:
        if not os.path.exists(f'results/{iter}'):
            os.mkdir(f'results/{iter}')
        denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054],
                                   std=[0.2726, 0.2778, 0.2861])
        img_id = 0
# =============================================================================
    criterion = nn.CrossEntropyLoss(ignore_index=255, reduction='mean')
# =============================================================================
    with torch.no_grad():
        torch.set_grad_enabled(True) 
        for i, (images, labels, filenames) in tqdm(enumerate(loader)):
            images = images.to(device, dtype=torch.float32)
            
#            images.requires_grad = True
         
#            torch.autograd.grad(images, create_graph=True, allow_unused=True)
            labels = labels.to(device, dtype=torch.long)
            new_images=Variable(images, requires_grad=True)
            
            new_labels=Variable(labels, requires_grad=False)

            outputs = model(new_images)
#            criterion = utils.FocalLoss(ignore_index=255, size_average=True)

            # torch.Size([4, 513, 513])    
            # according to the       
            #   st
            # mask = new_labels == 15
            # mask = mask.int()
            # np_mask = torch.unsqueeze(mask,1)
            #   ed
            # TODO how to get mask with the output?
            # print(outputs.shape)
            # print(new_labels.shape)
            # print(mask.shape)
            # print(mask[1,:,:][100])
            # plt.imshow(mask[1,:,:].cpu())
            # plt.show()
            # torch.Size([4, 21, 513, 513])

            # mask = outputs == 15
            # mask = mask.int()
            # mask = torch.max(mask,1).values
            # np_mask = torch.unsqueeze(mask,1)
            
            
            # print(np_mask.shape)
# =============================================================================
            # loss = criterion(outputs*(~np_mask), new_labels*(~mask)) + criterion(outputs*np_mask, new_labels*mask) 
            # loss_1 = criterion(outputs*(~ np_mask), new_labels*(~mask))
            #   st


            # loss_1 = criterion(outputs*(1 - np_mask), new_labels*(1 - mask))
            
            # print(((outputs * np_mask).shape))
            
            # print((new_labels.shape))

            # loss_2 = criterion(outputs * np_mask, new_labels*0)
            # t_loss = loss_1 + loss_2
            #   ed
            # print(t_loss)
            loss = criterion(outputs, new_labels)
               # Zero all existing gradients
            model.zero_grad()
   
               # Calculate gradients of model in backward pass
            # t_loss.backward()
            loss.backward()     
               # Collect datagrad
 ##            print(images.grad)
 #            sign_data_grad = torch.autograd.grad(loss, new_images,
 #                                       retain_graph=False, create_graph=False)[0]
 #            data_grad = new_images.grad.data
 #            sign_data_grad = torch.sign(data_grad)
   
               # Call FGSM Attack
               
            # TODO allocate the 1abel 15
            # TODO set up SegPGD
            # TODO turn to label not 15
            # TODO turn to label 0
            # adversarial_x = attacks.fgsm(images, new_images, 0.005)
            #
            # adversarial_x = attacks.t_fgsm_2(images, new_images, 4/255)
            adversarial_x = attacks.segpgd(images,new_images,new_labels,0.005,model)
            #
            # print(adversarial_x[1,1,:,:][2][100])
            # plt.imshow(adversarial_x[1,1,:,:].cpu())
            # plt.show()
            # adversarial_x = adversarial_x 
            # print(mask[1,:,:][100])
            # plt.imshow(mask[1,:,:].cpu())
            # plt.show()
            # print(np_mask.shape)
            # print(np_mask[0,0,:,:][200])
            # plt.imshow(np_mask[0,0,:,:].cpu())
            # plt.show()
#  new attack -> on the loss of the 


#            adversarial_x = attacks.pgd(images,new_images,new_labels,0.001,model)
               
               
 #            adversarial_x = images + 0.001 * sign_data_grad.sign_()
 #            adversarial_x = new_images + (0.005 * sign_data_grad)
             
 #            adversarial_y = new_images + 0.000000001 * sign_data_grad
     # Adding clipping to maintain [0,1] range
 #            adversarial_x = torch.clamp(adversarial_x, 0, 1)
 #            adversarial_y = torch.clamp(adversarial_y, 0, 1)
 #            adversarial_x = sign_data_grad
   
               # Re-classify the perturbed image
            new_output = model(adversarial_x)
             
            preds = new_output.detach().max(dim=1)[1].cpu().numpy()
# =============================================================================
#            preds = outputs.detach().max(dim=1)[1].cpu().numpy()
            targets = labels.cpu().numpy()

            metrics.update(targets, preds)
            
            if ret_samples_ids is not None and i in ret_samples_ids:  # get vis samples
                ret_samples.append(
                    (images[0].detach().cpu().numpy(), targets[0], preds[0]))

            if opts.save_val_results:
                for i in range(len(images)):
                    image = images[i].detach().cpu().numpy()
                    target = targets[i]
                    pred = preds[i]
# =============================================================================
                    adversarial_img = adversarial_x[i].detach().cpu().numpy()
# =============================================================================
                    
#                    adversarial_img_y =  adversarial_y[i].detach().cpu().numpy()

                    image = (denorm(image) * 255).transpose(1, 2, 0).astype(np.uint8)
# =============================================================================
                    adversarial_img = (denorm(adversarial_img) * 255).transpose(1, 2, 0).astype(np.uint8)
# =============================================================================
                    
#                    adversarial_img_y = (denorm(adversarial_img_y) * 255).transpose(1, 2, 0).astype(np.uint8)
                    
                    target = loader.dataset.decode_target(target).astype(np.uint8)
                    pred = loader.dataset.decode_target(pred).astype(np.uint8)

                    Image.fromarray(image).save(f'results/{iter}/{img_id}_image.png')
# =============================================================================
                    Image.fromarray(adversarial_img).save(f'results/{iter}/{img_id}_atimage.png')
# =============================================================================
#                    Image.fromarray(adversarial_img_y).save('results/%d_atyimage.png' % img_id)

                    Image.fromarray(target).save(f'results/{iter}/{img_id}_target.png')
                    Image.fromarray(pred).save(f'results/{iter}/{img_id}_pred.png')

                    fig = plt.figure()
                    plt.imshow(image)
                    plt.axis('off')
                    plt.imshow(pred, alpha=0.7)
                    ax = plt.gca()
                    ax.xaxis.set_major_locator(matplotlib.ticker.NullLocator())
                    ax.yaxis.set_major_locator(matplotlib.ticker.NullLocator())
                    plt.savefig(f'results/{iter}/{img_id}_overlay.png', bbox_inches='tight', pad_inches=0)
                    plt.close()
                    img_id += 1

        score = metrics.get_results()
    return score, ret_samples


def main():
    opts = get_argparser().parse_args()
    if opts.dataset.lower() == 'voc':
        opts.num_classes = 21
    elif opts.dataset.lower() == 'cityscapes':
        opts.num_classes = 19
    elif opts.dataset.lower() == 'trav':
        opts.num_classes = 19  # after loading the model, change the last MLP's dimension

    # Setup visualization
    vis = Visualizer(port=opts.vis_port,
                     env=opts.vis_env) if opts.enable_vis else None
    if vis is not None:  # display options
        vis.vis_table("Options", vars(opts))

    os.environ['CUDA_VISIBLE_DEVICES'] = opts.gpu_id
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    print(f"Device: {opts.gpu_id}")

    # Setup random seed
    torch.manual_seed(opts.random_seed)
    np.random.seed(opts.random_seed)
    random.seed(opts.random_seed)

    # Setup dataloader
    if opts.dataset == 'voc' and not opts.crop_val:
        opts.val_batch_size = 1

    train_dst, val_dst = get_dataset(opts)
    train_loader = data.DataLoader(
        train_dst, batch_size=opts.batch_size, shuffle=True, num_workers=2,
        drop_last=True)  # drop_last=True to ignore single-image batches.
    val_loader = data.DataLoader(
        val_dst, batch_size=opts.val_batch_size, shuffle=False, num_workers=2)
    print(f"Dataset: {opts.dataset}, Train set: {len(train_dst)}, Val set: {len(val_dst)}")

    # Set up model (all models are 'constructed at network.modeling)
    model = network.modeling.__dict__[opts.model](num_classes=opts.num_classes, output_stride=opts.output_stride)
    if opts.separable_conv and 'plus' in opts.model:
        network.convert_to_separable_conv(model.classifier)
    utils.set_bn_momentum(model.backbone, momentum=0.01)

    # Set up metrics
    metrics = StreamSegMetrics(opts.num_classes)

    # Set up optimizer
    optimizer = torch.optim.SGD(params=[
        {'params': model.backbone.parameters(), 'lr': 0.1 * opts.lr},
        {'params': model.classifier.parameters(), 'lr': opts.lr},
    ], lr=opts.lr, momentum=0.9, weight_decay=opts.weight_decay)
    # optimizer = torch.optim.SGD(params=model.parameters(), lr=opts.lr, momentum=0.9, weight_decay=opts.weight_decay)
    # torch.optim.lr_scheduler.StepLR(optimizer, step_size=opts.lr_decay_step, gamma=opts.lr_decay_factor)
    if opts.lr_policy == 'poly':
        scheduler = utils.PolyLR(optimizer, opts.total_itrs, power=0.9)
    elif opts.lr_policy == 'step':
        scheduler = torch.optim.lr_scheduler.StepLR(optimizer, step_size=opts.step_size, gamma=0.1)

    # Set up criterion
    # criterion = utils.get_loss(opts.loss_type)
    if opts.loss_type == 'focal_loss':
        criterion = utils.FocalLoss(ignore_index=255, size_average=True)
    elif opts.loss_type == 'cross_entropy':
        criterion = nn.CrossEntropyLoss(ignore_index=255, reduction='mean')

    def save_ckpt(path):
        """ save current model
        """
        torch.save({
            "cur_itrs": cur_itrs,
            "model_state": model.module.state_dict(),
            "optimizer_state": optimizer.state_dict(),
            "scheduler_state": scheduler.state_dict(),
            "best_score": best_score,
        }, path)
        print("Model saved as %s" % path)

    utils.mkdir('checkpoints')
    # Restore
    best_score = 0.0
    cur_itrs = 0
    cur_epochs = 0
    if opts.ckpt is not None and os.path.isfile(opts.ckpt):
        # https://github.com/VainF/DeepLabV3Plus-Pytorch/issues/8#issuecomment-605601402, @PytaichukBohdan
        checkpoint = torch.load(opts.ckpt, map_location=torch.device('cpu'))
        model.load_state_dict(checkpoint["model_state"])
        model.classifier.classifier[-1] = nn.Conv2d(256, 2, 1)
        model = nn.DataParallel(model)
        model.to(device)

        if opts.continue_training:
            optimizer.load_state_dict(checkpoint["optimizer_state"])
            scheduler.load_state_dict(checkpoint["scheduler_state"])
            cur_itrs = checkpoint["cur_itrs"]
            best_score = checkpoint['best_score']
            print("Training state restored from %s" % opts.ckpt)
        print("Model restored from %s" % opts.ckpt)
        del checkpoint  # free memory
    else:
        print("[!] Retrain")
        model = nn.DataParallel(model)
        model.to(device)

    # ==========   Train Loop   ==========#
    vis_sample_id = np.random.randint(0, len(val_loader), opts.vis_num_samples,
                                      np.int32) if opts.enable_vis else None  # sample idxs for visualization
    denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054], std=[0.2726, 0.2778, 0.2861])  # denormalization for ori images

    if opts.test_only:
        model.eval()
        val_score, ret_samples = validate(
            opts, model, val_loader, device, metrics, cur_itrs, vis_sample_id)
        print(metrics.to_str(val_score))
        return

    interval_loss = 0
    # https://github.com/ndb796/Pytorch-Adversarial-Training-CIFAR/blob/master/pgd_adversarial_training.py
    while True:  # cur_itrs < opts.total_itrs:
        # =====  Train  =====
        model.train()
        cur_epochs += 1
# =============================================================================
#         #s
#         torch.set_grad_enabled(True) 
#         #e
# =============================================================================
        for (images, labels, filenames) in train_loader:
            cur_itrs += 1

            images = images.to(device, dtype=torch.float32)
            labels = labels.to(device, dtype=torch.long)
            
# =============================================================================
#             # s
#             new_images=Variable(images, requires_grad=True)
#             new_labels=Variable(labels, requires_grad=False)
#             # e
# =============================================================================
            
            optimizer.zero_grad()
#            outputs = model(new_images)
            
            # s
            
# =============================================================================
#             #
#             new_images_d = new_images.detach()
#             new_images_d.requires_grad_()
#             with torch.enable_grad():
#                 outputs_a = model(new_images_d)
#                 loss_a = criterion(outputs_a, labels)
#             data_grad = torch.autograd.grad(loss_a, [new_images_d])[0]
#             adversarial_x = new_images_d.detach() + 0.005 * torch.sign(data_grad.detach())
#             new_output = model(adversarial_x)
#             #
# =============================================================================
            
            # e
            
#            loss = criterion(new_output, labels)
            
            outputs = model(images)
            
#            loss_o.backward()
            
            # s
            
#            data_grad = torch.autograd.grad(loss_o, [new_images])[0]
#            adversarial_x = attacks.fgsm(images, data_grad, 0.005)
#            new_output = model(adversarial_x)
# =============================================================================
#             #
#             lamb = 0.5
#             loss = (1-lamb) * criterion(outputs, labels) + lamb * criterion(new_output, labels)
#             #
# =============================================================================

            loss = criterion(outputs[0], labels)
            loss.backward()

            optimizer.step()

            np_loss = loss.detach().cpu().numpy()
            interval_loss += np_loss
            if vis is not None:
                vis.vis_scalar('Loss', cur_itrs, np_loss)

            if (cur_itrs) % 10 == 0:
                interval_loss = interval_loss / 10
                print(f"Epoch {cur_epochs}, Itrs {cur_itrs}/{opts.total_itrs}, Loss={interval_loss}")
                interval_loss = 0.0

            if (cur_itrs) % opts.val_interval == 0:
                save_ckpt(f'checkpoints/latest_{opts.model}_{opts.dataset}_os{opts.output_stride}.pth')
                print("validation...")
                model.eval()
                val_score, ret_samples = validate(
                    opts, model, val_loader, device, metrics, cur_itrs,
                    vis_sample_id)
                print(metrics.to_str(val_score))
                if val_score['Mean IoU'] > best_score:  # save best model
                    best_score = val_score['Mean IoU']
                    save_ckpt('checkpoints/best_%s_%s_os%d.pth' %
                              (opts.model, opts.dataset, opts.output_stride))

                if vis is not None:  # visualize validation score and samples
                    vis.vis_scalar("[Val] Overall Acc", cur_itrs, val_score['Overall Acc'])
                    vis.vis_scalar("[Val] Mean IoU", cur_itrs, val_score['Mean IoU'])
                    vis.vis_table("[Val] Class IoU", val_score['Class IoU'])

                    for k, (img, target, lbl) in enumerate(ret_samples):
                        img = (denorm(img) * 255).astype(np.uint8)
                        target = train_dst.decode_target(target).transpose(2, 0, 1).astype(np.uint8)
                        lbl = train_dst.decode_target(lbl).transpose(2, 0, 1).astype(np.uint8)
                        concat_img = np.concatenate((img, target, lbl), axis=2)  # concat along width
                        vis.vis_image('Sample %d' % k, concat_img)
                model.train()
            scheduler.step()

            if cur_itrs >= opts.total_itrs:
                save_ckpt(f'checkpoints/last_{opts.model}_{opts.dataset}_os{opts.output_stride}.pth')
                return


def per_image_metric(metric, y_true, y_pred):
    """
    calculate some metrics for one image
    """
    metric.reset()
    hist = metric._fast_hist(y_true.flatten(), y_pred.flatten())  # confusion matrix
    acc = np.diag(hist).sum() / hist.sum()
    acc_cls = np.diag(hist) / hist.sum(axis=1)
    acc_cls = np.nanmean(acc_cls)
    iu = np.diag(hist) / (hist.sum(axis=1) + hist.sum(axis=0) - np.diag(hist))
    cls_iu = dict(zip(range(metric.n_classes), iu))
    mean_iu = np.nanmean(iu)
    return {
            "Overall Acc": acc,
            "Mean Acc": acc_cls,
            "Mean IoU": mean_iu,
            **cls_iu
            }


def inference(split='val'):
    """
    load saved pretrained model
    infer on train/val sets
    """
    opts = get_argparser().parse_args()
    if opts.dataset.lower() == 'trav':
        opts.num_classes = 2
    
    os.environ['CUDA_VISIBLE_DEVICES'] = opts.gpu_id
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    print(f"Device: {opts.gpu_id}")

    train_dst, val_dst = get_dataset(opts)
    if split == 'train':
        loader = data.DataLoader(
            train_dst, batch_size=opts.batch_size, shuffle=True, num_workers=2,
            drop_last=True)  # drop_last=True to ignore single-image batches.
    else:
        loader = data.DataLoader(
            val_dst, batch_size=opts.val_batch_size, shuffle=False, num_workers=2)
    print(f"Dataset: {opts.dataset}, Train set: {len(train_dst)}, Val set: {len(val_dst)}")

    metrics = StreamSegMetrics(opts.num_classes)
    perimage_metrics = StreamSegMetrics(opts.num_classes)

    model = network.modeling.__dict__[opts.model](num_classes=opts.num_classes, output_stride=opts.output_stride)
    utils.set_bn_momentum(model.backbone, momentum=0.01)
    checkpoint = torch.load(f'checkpoints/hloss_{opts.model}_{opts.dataset}_os{opts.output_stride}.pt', map_location=torch.device('cpu'))
    model.load_state_dict(checkpoint["model_state"])
    # model.classifier.classifier[-1] = nn.Conv2d(256, 2, 1)
    model = nn.DataParallel(model)
    model.to(device)
    model.eval()

    vis_sample_id = np.random.randint(0, len(loader), opts.vis_num_samples,
                                      np.int32) if opts.enable_vis else None
    # from validate fn
    metrics.reset()
    ret_samples = []
    if opts.save_val_results:
        if not os.path.exists(f'results/{opts.dataset}'):
            os.mkdir(f'results/{opts.dataset}')
        denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054],
                                   std=[0.2726, 0.2778, 0.2861])
        img_id = 0

    # criterion = nn.CrossEntropyLoss(ignore_index=255, reduction='mean')
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    left = top = Inches(0.1)
    table_top = Inches(6)
    width = Inches(14.0)
    height = Inches(1.2)
    # df = pd.DataFrame(columns=['image_id', 'eps', 'alpha', 'effect'])

    for i, (x_clean, y_true, filenames) in tqdm(enumerate(loader)):
        x_clean = x_clean.to(device, dtype=torch.float32)
        y_true = y_true.to(device, dtype=torch.uint8)
        
        y_pred_clean = model(x_clean)

        delta1 = attacks.pgd(model, x_clean, y_true, device, epsilon=opts.eps, alpha=opts.alpha, num_iter=10)
        x_adv = x_clean.float() + delta1.float()
        y_pred_adv = model(x_adv)
        y_pred_adv_np = y_pred_adv.detach().max(dim=1)[1].cpu().numpy()
        y_true_np = y_true.detach().cpu().numpy()
        y_pred_clean_np = y_pred_clean.detach().max(dim=1)[1].cpu().numpy()
        metrics.update(y_true_np, y_pred_adv_np)

        if vis_sample_id is not None and i in vis_sample_id:  # get vis samples
            ret_samples.append(
                (x_clean[0].detach().cpu().numpy(), y_true_np[0], y_pred_adv_np[0]))

        if opts.save_val_results:
            for j in range(len(x_clean)):
                image = x_clean[j].detach().cpu().numpy()
                delta_np = np.sum(delta1[j].detach().cpu().numpy(), axis=0)
                total_pertub = np.sum(np.abs(delta_np))
                target = y_true_np[j]  # y_true
                pred = y_pred_adv_np[j]  # adv y_pred
                output = y_pred_clean_np[j]  # clean y_pred
                adv_iou = per_image_metric(perimage_metrics, target, pred)
                clean_iou = per_image_metric(perimage_metrics, target, output)
                adversarial_img = x_adv[j].detach().cpu().numpy()
                effect = (clean_iou['Mean IoU'] - adv_iou['Mean IoU'])/total_pertub

                image = (denorm(image) * 255).transpose(1, 2, 0).astype(np.uint8)

                adversarial_img = (denorm(adversarial_img) * 255).transpose(1, 2, 0).astype(np.uint8)

                slide = prs.slides.add_slide(blank_slide_layout)
                fig, axs = plt.subplots(2, 3, figsize=(14, 6))
                axs[0,0].imshow(image)
                axs[0,0].set_title(f'Clean image')
                axs[0,0].axis('off')

                axs[0,1].imshow(image)
                axs[0,1].imshow(target, cmap='viridis', alpha=0.4)
                axs[0,1].set_title(f'y_true')
                axs[0,1].axis('off')

                axs[0,2].imshow(image)
                axs[0,2].imshow(output, cmap='viridis', alpha=0.4)
                axs[0,2].set_title(f'Clean y_pred')
                axs[0,2].axis('off')

                axs[1,0].imshow(adversarial_img)
                axs[1,0].set_title(f'Adv image')
                axs[1,0].axis('off')

                axs[1,1].imshow(adversarial_img)
                axs[1,1].imshow(pred, cmap='viridis', alpha=0.4)
                axs[1,1].set_title(f'Adv y_pred')
                axs[1,1].axis('off')

                axs[1,2].set_title(f'delta, eps={opts.eps}, alpha={opts.alpha}, total_pertub: {total_pertub}')
                axs[1,2].imshow(delta_np, cmap='viridis')
                axs[1,2].axis('off')

                img_filename = f'results/{opts.dataset}/{img_id}_overlay.png'
                fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
                plt.close()
                file_parts = filenames[j].split(os.sep)
                split_index = file_parts.index('segmentation_indoor_images')
                right_filename = os.sep.join(file_parts[split_index+1:])
                pic = slide.shapes.add_picture(img_filename, left, top)
                shapes = slide.shapes
                table = shapes.add_table(4, 6, left, table_top, width, height).table
                table.cell(0, 0).text = right_filename
                table.cell(1, 0).text = 'Clean'
                table.cell(2, 0).text = 'Adv'
                table.cell(3, 0).text = 'total delta mean IoU/pertub'
                table.cell(3, 1).text = f"{effect}"
                keys = list(clean_iou.keys())
                for idx, k in enumerate(keys):
                    table.cell(0, idx+1).text = str(k)  # [Overall acc, ...]
                    table.cell(1, idx+1).text = f'{clean_iou[k]:.4f}'
                    table.cell(2, idx+1).text = f'{adv_iou[k]:.4f}'

                # append to df
                # new_df = pd.DataFrame({'image_id': right_filename, 'eps': eps, 'alpha': alfa, 'effect': effect}, index=[1])
                # df = pd.concat([df, new_df], ignore_index=True)
                img_id += 1

        # if i > 10:
        #     break

    score = metrics.get_results()
    print(score)
    # df.to_csv(f'results/effects.csv')
    prs.save(f'results/{opts.dataset}_{split}6.pptx')
    return score, ret_samples


def save_separate_images(split='val'):
    """
    similar to the above fn, but save separate images instead of together.
    """
    model_type = f'hloss'  # or clean
    opts = get_argparser().parse_args()
    if opts.dataset.lower() == 'trav':
        opts.num_classes = 2
    
    os.environ['CUDA_VISIBLE_DEVICES'] = opts.gpu_id
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    print(f"Device: {opts.gpu_id}")

    train_dst, val_dst = get_dataset(opts)
    if split == 'train':
        loader = data.DataLoader(
            train_dst, batch_size=opts.batch_size, shuffle=True, num_workers=2,
            drop_last=True)  # drop_last=True to ignore single-image batches.
    else:
        loader = data.DataLoader(
            val_dst, batch_size=opts.val_batch_size, shuffle=False, num_workers=2)
    print(f"Dataset: {opts.dataset}, Train set: {len(train_dst)}, Val set: {len(val_dst)}")

    model = network.modeling.__dict__[opts.model](num_classes=opts.num_classes, output_stride=opts.output_stride)
    utils.set_bn_momentum(model.backbone, momentum=0.01)
    checkpoint = torch.load(f'checkpoints/{model_type}_{opts.model}_{opts.dataset}_os{opts.output_stride}.pt', map_location=torch.device('cpu'))
    model.load_state_dict(checkpoint["model_state"])
    model = nn.DataParallel(model)
    model.to(device)
    model.eval()

    # from validate fn
    if opts.save_val_results:
        if not os.path.exists(f'results/{model_type}'):
            os.mkdir(f'results/{model_type}')
        denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054],
                                   std=[0.2726, 0.2778, 0.2861])
        img_id = 0

    for i, (x_clean, y_true, filenames) in tqdm(enumerate(loader), desc=f'batches'):
        x_clean = x_clean.to(device, dtype=torch.float32)
        y_true = y_true.to(device, dtype=torch.uint8)

        y_pred_clean, _ = model(x_clean)

        delta1 = attacks.pgd(model, x_clean, y_true, device, epsilon=opts.eps, alpha=opts.alpha, num_iter=10)
        x_adv = x_clean.float() + delta1.float()
        y_pred_adv, _ = model(x_adv)

        y_pred_adv_np = y_pred_adv.detach().max(dim=1)[1].cpu().numpy()
        y_true_np = y_true.detach().cpu().numpy()
        y_pred_clean_np = y_pred_clean.detach().max(dim=1)[1].cpu().numpy()

        if opts.save_val_results:
            for j in range(len(x_clean)):
                image = x_clean[j].detach().cpu().numpy()
                # delta_np = np.sum(delta1[j].detach().cpu().numpy(), axis=0)
                # total_pertub = np.sum(np.abs(delta_np))
                target = y_true_np[j]  # y_true
                pred = y_pred_adv_np[j]  # adv y_pred
                output = y_pred_clean_np[j]  # clean y_pred

                adversarial_img = x_adv[j].detach().cpu().numpy()
                image = (denorm(image) * 255).transpose(1, 2, 0).astype(np.uint8)
                adversarial_img = (denorm(adversarial_img) * 255).transpose(1, 2, 0).astype(np.uint8)

                # plt.imsave(f'results/trav/{model_type}_{i}_{j}_x_clean.png', image)

                # fig, ax = plt.subplots()
                # ax.imshow(image)
                # ax.imshow(target, cmap='viridis', alpha=0.4)
                # ax.axis('off')
                # fig.savefig(f'results/trav/{i}_{j}_y_true.png',bbox_inches='tight', pad_inches=0)
                # plt.close(fig)

                fig, ax = plt.subplots()
                ax.imshow(image)
                ax.imshow(output, cmap='viridis', alpha=0.4)
                ax.axis('off')
                fig.savefig(f'results/{model_type}/{model_type}_{i}_{j}_clean_y_pred.png',bbox_inches='tight', pad_inches=0)
                plt.close(fig)

                # plt.imsave(f'results/trav/{i}_{j}_x_adv.png', adversarial_img)

                fig, ax = plt.subplots()
                ax.imshow(adversarial_img)
                ax.imshow(pred, cmap='viridis', alpha=0.4)
                ax.axis('off')
                fig.savefig(f'results/{model_type}/{model_type}_{i}_{j}_adv_y_pred.png',bbox_inches='tight', pad_inches=0)
                plt.close(fig)

                img_id += 1

    return


def increasing_perturbations(split='val'):
    """
    generate perturbed images with eps=[0.005, 0.05, 0.1, 0.2, 0.3] and alpha=100
    """
    model_type = f'clean'  # or clean
    opts = get_argparser().parse_args()
    if opts.dataset.lower() == 'trav':
        opts.num_classes = 2
    
    os.environ['CUDA_VISIBLE_DEVICES'] = opts.gpu_id
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    print(f"Device: {opts.gpu_id}")

    train_dst, val_dst = get_dataset(opts)
    if split == 'train':
        loader = data.DataLoader(
            train_dst, batch_size=opts.batch_size, shuffle=True, num_workers=2,
            drop_last=True)  # drop_last=True to ignore single-image batches.
    else:
        loader = data.DataLoader(
            val_dst, batch_size=opts.val_batch_size, shuffle=False, num_workers=2)
    print(f"Dataset: {opts.dataset}, Train set: {len(train_dst)}, Val set: {len(val_dst)}")

    model = network.modeling.__dict__[opts.model](num_classes=opts.num_classes, output_stride=opts.output_stride)
    utils.set_bn_momentum(model.backbone, momentum=0.01)
    checkpoint = torch.load(f'checkpoints/{model_type}_{opts.model}_{opts.dataset}_os{opts.output_stride}.pt', map_location=torch.device('cpu'))
    model.load_state_dict(checkpoint["model_state"])
    model = nn.DataParallel(model)
    model.to(device)
    model.eval()

    # from validate fn
    if opts.save_val_results:
        if not os.path.exists(f'results/{model_type}'):
            os.mkdir(f'results/{model_type}')
        denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054],
                                   std=[0.2726, 0.2778, 0.2861])
        img_id = 0

    epsilons = [0.005,0.01, 0.05, 0.1]

    for i, (x_clean, y_true, filenames) in tqdm(enumerate(loader), desc=f'batches'):
        if i == 107:
            x_clean = x_clean.to(device, dtype=torch.float32)
            y_true = y_true.to(device, dtype=torch.uint8)

            # y_pred_clean, clean_features = model(x_clean)
            for eps in epsilons:
                delta1 = attacks.pgd(model, x_clean, y_true, device, epsilon=eps, alpha=opts.alpha, num_iter=10)
                x_adv = x_clean.float() + delta1.float()
                # y_pred_adv, adv_features = model(x_adv)

                # y_pred_adv_np = y_pred_adv.detach().max(dim=1)[1].cpu().numpy()
                # y_true_np = y_true.detach().cpu().numpy()
                # y_pred_clean_np = y_pred_clean.detach().max(dim=1)[1].cpu().numpy()

                if opts.save_val_results:
                    for j in range(len(x_clean)):
                        image = x_clean[j].detach().cpu().numpy()

                        adversarial_img = x_adv[j].detach().cpu().numpy()
                        image = (denorm(image) * 255).transpose(1, 2, 0).astype(np.uint8)
                        adversarial_img = (denorm(adversarial_img) * 255).transpose(1, 2, 0).astype(np.uint8)

                        plt.imsave(f'results/{model_type}/{i}_{j}_x_adv_{eps}.png', adversarial_img)

                        img_id += 1

    return


def save_qualitative_results(split='val'):
    """
    3 positives: 80_2, 58_1, 92_3
    1 negative: 75_1
    save separate images: clean image, adv image, perturbation, label
    (before defense, after AT, after AT+hloss) existed, no need
    """
    model_type = f'hloss'  # or clean
    opts = get_argparser().parse_args()
    if opts.dataset.lower() == 'trav':
        opts.num_classes = 2
    
    os.environ['CUDA_VISIBLE_DEVICES'] = opts.gpu_id
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    print(f"Device: {opts.gpu_id}")

    train_dst, val_dst = get_dataset(opts)
    if split == 'train':
        loader = data.DataLoader(
            train_dst, batch_size=opts.batch_size, shuffle=True, num_workers=2,
            drop_last=True)  # drop_last=True to ignore single-image batches.
    else:
        loader = data.DataLoader(
            val_dst, batch_size=opts.val_batch_size, shuffle=False, num_workers=2)
    print(f"Dataset: {opts.dataset}, Train set: {len(train_dst)}, Val set: {len(val_dst)}")

    model = network.modeling.__dict__[opts.model](num_classes=opts.num_classes, output_stride=opts.output_stride)
    utils.set_bn_momentum(model.backbone, momentum=0.01)
    checkpoint = torch.load(f'checkpoints/{model_type}_{opts.model}_{opts.dataset}_os{opts.output_stride}.pt', map_location=torch.device('cpu'))
    model.load_state_dict(checkpoint["model_state"])
    model = nn.DataParallel(model)
    model.to(device)
    model.eval()

    # from validate fn
    if opts.save_val_results:
        if not os.path.exists(f'results/{model_type}'):
            os.mkdir(f'results/{model_type}')
        denorm = utils.Denormalize(mean=[0.5174, 0.4857, 0.5054],
                                   std=[0.2726, 0.2778, 0.2861])
        img_id = 0

    # selected_images = ['0_0']  # '80_2', '58_1', '91_3', '75_1'
    selected_images = ['7_0','14_2','39_2','79_0','86_1']
    fig, axs = plt.subplots(8, 4)  # figsize=(15, 6)
    for i, (x_clean, y_true, filenames) in tqdm(enumerate(loader), desc=f'batches'):
        x_clean = x_clean.to(device, dtype=torch.float32)
        y_true = y_true.to(device, dtype=torch.uint8)

        y_pred_clean, _ = model(x_clean)

        delta1 = attacks.pgd(model, x_clean, y_true, device, epsilon=opts.eps, alpha=opts.alpha, num_iter=20)
        x_adv = x_clean.float() + delta1.float()
        y_pred_adv, _ = model(x_adv)

        y_pred_adv_np = y_pred_adv.detach().max(dim=1)[1].cpu().numpy()
        y_true_np = y_true.detach().cpu().numpy()
        y_pred_clean_np = y_pred_clean.detach().max(dim=1)[1].cpu().numpy()

        if opts.save_val_results:
            for j in range(len(x_clean)):
                if f'{i}_{j}' in selected_images:
                    image = x_clean[j].detach().cpu().numpy()
                    delta_np = np.sum(np.abs(delta1[j].detach().cpu().numpy()), axis=0)
                    target = y_true_np[j]  # y_true
                    pred = y_pred_adv_np[j]  # adv y_pred
                    output = y_pred_clean_np[j]  # clean y_pred

                    adversarial_img = x_adv[j].detach().cpu().numpy()
                    image = (denorm(image) * 255).transpose(1, 2, 0).astype(np.uint8)
                    adversarial_img = (denorm(adversarial_img) * 255).transpose(1, 2, 0).astype(np.uint8)

                    plt.imsave(f'results/qualitative_results/{i}_{j}_x_clean.png', image)

                    plt.imsave(f'results/qualitative_results/{i}_{j}_x_adv.png', adversarial_img)

                    fig, ax = plt.subplots()
                    ax.imshow(image)
                    ax.imshow(target, cmap='viridis', alpha=0.4)
                    ax.axis('off')
                    fig.savefig(f'results/qualitative_results/{i}_{j}_y_true.png',bbox_inches='tight', pad_inches=0)
                    plt.close(fig)

                    fig, ax = plt.subplots()
                    ax.imshow(delta_np, cmap='viridis')
                    ax.axis('off')
                    fig.savefig(f'results/qualitative_results/{i}_{j}_delta.png',bbox_inches='tight', pad_inches=0)
                    plt.close(fig)


if __name__ == '__main__':
    main()
    # inference()
    # save_separate_images()
    # increasing_perturbations()
    # save_qualitative_results()
